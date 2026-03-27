"""Generate ProjectQ one-pager slide from PDF content using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x12, 0x1E, 0x3C)
TEAL       = RGBColor(0x02, 0xA8, 0x8A)
ORANGE     = RGBColor(0xF0, 0x7A, 0x30)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB8, 0xC4, 0xD4)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
BODY_TEXT  = RGBColor(0x44, 0x44, 0x58)
STAT_TEAL  = RGBColor(0x00, 0xB4, 0x96)
BLUE_DOT   = RGBColor(0x02, 0x8F, 0xBF)
PURPLE     = RGBColor(0x7C, 0x5C, 0xBF)
DIVIDER    = RGBColor(0xDC, 0xE4, 0xEE)
MUTED_NAVY = RGBColor(0x60, 0x6A, 0x80)
TEAM_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
IMPACT_BG  = RGBColor(0xF0, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(x, y, w, h, fill=None, line_color=None, line_width=Pt(0.5)):
    sp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = line_width
    else:
        sp.line.fill.background()
    return sp


def new_tf(x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    return tf


def run(p, text, size=9, bold=False, italic=False, color=DARK_TEXT):
    r = p.add_run()
    r.text = text
    r.font.size     = Pt(size)
    r.font.bold     = bold
    r.font.italic   = italic
    r.font.color.rgb = color
    return r


def simple(x, y, w, h, text, size=10, bold=False, italic=False,
           color=DARK_TEXT, align=PP_ALIGN.LEFT):
    tf = new_tf(x, y, w, h)
    p  = tf.paragraphs[0]
    p.alignment = align
    run(p, text, size=size, bold=bold, italic=italic, color=color)
    return tf


def para(tf, text, size=9, bold=False, italic=False,
         color=BODY_TEXT, align=PP_ALIGN.LEFT, space_before=3):
    p = tf.add_paragraph()
    p.alignment    = align
    p.space_before = Pt(space_before)
    run(p, text, size=size, bold=bold, italic=italic, color=color)
    return p


# ── Background ────────────────────────────────────────────────────────────────
add_rect(0, 0, 13.33, 7.5, fill=WHITE)

# ── Header ────────────────────────────────────────────────────────────────────
add_rect(0, 0, 13.33, 1.35, fill=NAVY)

tf = new_tf(0.4, 0.10, 8.5, 0.72)
p1 = tf.paragraphs[0]
run(p1, "ProjectQ", size=30, bold=True, color=WHITE)
p2 = tf.add_paragraph()
p2.space_before = Pt(3)
run(p2, "The invisible member of your engineering team", size=11, color=LIGHT_GRAY)

simple(0.4, 0.90, 9.2, 0.24,
       "Platform engineering meets AI. One capability. Every team gets an invisible +1. Sprint value goes up.",
       size=8.5, italic=True, color=RGBColor(0x80, 0x8A, 0xA0))

add_rect(0.4, 1.22, 1.0, 0.055, fill=TEAL)  # accent line

# N+1 badge
add_rect(11.35, 0.13, 1.6, 0.9, fill=ORANGE)
simple(11.35, 0.24, 1.6, 0.68, "N + 1", size=26, bold=True, color=WHITE,
       align=PP_ALIGN.CENTER)

# ── Column divider ────────────────────────────────────────────────────────────
add_rect(5.93, 1.38, 0.025, 4.48, fill=DIVIDER)

# ── LEFT COLUMN ───────────────────────────────────────────────────────────────
LX = 0.4
LW = 5.38
CY = 1.45

# Section: WHY PLATFORM ENGINEERING
simple(LX, CY, LW, 0.22, "WHY PLATFORM ENGINEERING", size=9, bold=True, color=DARK_TEXT)
add_rect(LX, CY + 0.23, 1.55, 0.04, fill=ORANGE)

# Pain points (bold label + regular text in same paragraph)
tf = new_tf(LX, CY + 0.35, LW, 1.55)
pain_points = [
    ("Cognitive Overload",
     " — Developers waste energy choosing, configuring, and maintaining AI tools instead of solving business problems."),
    ("No Integrated Ecosystem",
     " — Teams have isolated AI tools, not a connected experience. No shared context, no continuity between refinement, coding, review, and delivery."),
    ("Sprint Value is Hard to Protect",
     " — Maintaining alignment between what's planned and delivered requires constant attention across multiple tools and handoffs."),
    ("Every Team Starts from Zero",
     " — Without platform-provided AI, each team builds their own workflow. Duplication, inconsistency, no standards."),
]
for i, (bold_part, reg_part) in enumerate(pain_points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    if i > 0:
        p.space_before = Pt(3)
    run(p, bold_part, size=8, bold=True, color=DARK_TEXT)
    run(p, reg_part, size=8, color=BODY_TEXT)

# Section: THE IDEA
IDEA_Y = CY + 2.05
simple(LX, IDEA_Y, LW, 0.22, "THE IDEA", size=9, bold=True, color=DARK_TEXT)
add_rect(LX, IDEA_Y + 0.23, 0.65, 0.04, fill=TEAL)

simple(LX, IDEA_Y + 0.33, LW, 0.78,
       "ProjectQ is the kiss-and-ride experience for AI-native development. "
       "Platform engineering builds it once. Teams step in and start delivering. "
       "Minimal setup. No cognitive load. Just an invisible +1 member that watches "
       "over story quality, code, tests, security, and sprint outcomes — so developers "
       "can focus entirely on what matters: business value and quality.",
       size=8, color=BODY_TEXT)

# Team diagram
TEAM_Y = IDEA_Y + 1.18
for i, role in enumerate(["Dev", "Dev", "QA", "PO", "SM"]):
    rx = LX + i * 0.86
    add_rect(rx, TEAM_Y, 0.68, 0.30, fill=TEAM_GRAY)
    simple(rx, TEAM_Y + 0.05, 0.68, 0.20, role, size=7.5,
           color=DARK_TEXT, align=PP_ALIGN.CENTER)

plus_x = LX + 5 * 0.86 + 0.05
simple(plus_x, TEAM_Y, 0.28, 0.30, "+", size=13, bold=True,
       color=TEAL, align=PP_ALIGN.CENTER)
qx = plus_x + 0.32
add_rect(qx, TEAM_Y, 0.68, 0.30, fill=TEAL)
simple(qx, TEAM_Y + 0.05, 0.68, 0.20, "ProjectQ", size=7, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

# ── RIGHT COLUMN ──────────────────────────────────────────────────────────────
RX   = 6.08
RW   = 6.88

simple(RX, CY, RW, 0.22, "SPRINT AGENT MAP", size=9, bold=True, color=DARK_TEXT)
add_rect(RX, CY + 0.23, 1.35, 0.04, fill=ORANGE)

phases = [
    {
        "title": "PRE-SPRINT", "color": TEAL,
        "agents": [
            ("Story Quality Analyst", "Enriches ACs, flags edge cases"),
            ("Dependency Scanner", "Identifies conflicts before sprint"),
            ("Sprint Goal Synthesizer", "Summarizes business value"),
        ],
    },
    {
        "title": "DURING SPRINT", "color": BLUE_DOT,
        "agents": [
            ("Code Companion", "Suggests patterns from codebase"),
            ("MR Review Agent", "Reviews against story intent + standards"),
            ("Test Generator", "AC-aligned tests on commit"),
            ("Security Sentinel", "Daily scan, auto-fix MRs"),
            ("Standup Prep", "Summaries from Git + Jira activity"),
        ],
    },
    {
        "title": "POST-SPRINT", "color": PURPLE,
        "agents": [
            ("AC Validator", "Verifies 'done' matches criteria"),
            ("Demo Deck Generator", "Stakeholder-ready presentation"),
            ("Retro Insights", "Data-driven sprint patterns"),
        ],
    },
]

COL_W   = 2.1
COL_GAP = 0.19
PHASE_Y = CY + 0.38

for i, phase in enumerate(phases):
    cx = RX + i * (COL_W + COL_GAP)
    simple(cx, PHASE_Y, COL_W, 0.20, phase["title"],
           size=7.5, bold=True, color=phase["color"])
    add_rect(cx, PHASE_Y + 0.22, COL_W, 0.025, fill=phase["color"])

    ay = PHASE_Y + 0.30
    for name, desc in phase["agents"]:
        add_rect(cx, ay + 0.04, 0.10, 0.10, fill=phase["color"])
        tf = new_tf(cx + 0.17, ay, COL_W - 0.18, 0.36)
        p  = tf.paragraphs[0]
        run(p, name, size=7.5, bold=True, color=DARK_TEXT)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        run(p2, desc, size=7, color=BODY_TEXT)
        ay += 0.42

# ── PROJECTED IMPACT ──────────────────────────────────────────────────────────
IMP_Y = 5.88
add_rect(0, IMP_Y, 13.33, 1.02, fill=IMPACT_BG)

simple(0.4, IMP_Y + 0.08, 2.2, 0.22, "PROJECTED IMPACT",
       size=8, bold=True, color=DARK_TEXT)
add_rect(0.4, IMP_Y + 0.30, 1.2, 0.04, fill=ORANGE)

stats = [
    ("30%",    ["Faster sprint", "planning"]),
    ("50%",    ["Faster first", "MR review"]),
    ("85%+",   ["Test coverage", "on new code"]),
    ("<20min", ["Demo prep", "(vs 2+ hours)"]),
    ("40%",    ["Faster dev", "onboarding"]),
]

S_START = 2.7
S_GAP   = 2.15
for i, (num, label_lines) in enumerate(stats):
    sx = S_START + i * S_GAP
    simple(sx, IMP_Y + 0.06, 1.9, 0.42, num,
           size=22, bold=True, color=STAT_TEAL, align=PP_ALIGN.CENTER)
    tf = new_tf(sx, IMP_Y + 0.54, 1.9, 0.40)
    for j, line in enumerate(label_lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, size=7.5, color=BODY_TEXT)

# ── FOOTER ────────────────────────────────────────────────────────────────────
FOOT_Y = 6.93
add_rect(0, FOOT_Y, 13.33, 0.57, fill=NAVY)

simple(0.5, FOOT_Y + 0.07, 12.3, 0.28,
       '"Platform engineering builds the paths teams walk on. ProjectQ makes those paths '
       'AI-native. Your team is N — ProjectQ makes it N+1."',
       size=8, italic=True, color=RGBColor(0x90, 0x9A, 0xB0), align=PP_ALIGN.CENTER)

simple(0.5, FOOT_Y + 0.33, 5.5, 0.18,
       "ProjectQ  |  Innovation Day 2026  |  Confidential",
       size=6.5, color=MUTED_NAVY)

simple(7.33, FOOT_Y + 0.33, 5.5, 0.18, "Platform Engineering",
       size=6.5, color=MUTED_NAVY, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
output = "ProjectQ_OnePager.pptx"
prs.save(output)
print(f"Saved: {output}")
